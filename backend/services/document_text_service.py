from pathlib import Path

import fitz
from docx import Document as DocxDocument
from pptx import Presentation


SUPPORTED_DOCUMENT_TYPES = {
    ".pdf",
    ".docx",
    ".pptx",
}


# ==================================================
# EXTRACT PDF TEXT
# ==================================================

def extract_pdf_text(
    file_path: Path,
) -> str:
    text_parts = []

    with fitz.open(file_path) as pdf:
        for page in pdf:
            page_text = page.get_text("text")

            if page_text.strip():
                text_parts.append(page_text)

    return "\n\n".join(text_parts)


# ==================================================
# EXTRACT DOCX TEXT
# ==================================================

def extract_docx_text(
    file_path: Path,
) -> str:
    document = DocxDocument(file_path)

    text_parts = []

    # Paragraphs
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()

        if text:
            text_parts.append(text)

    # Tables
    for table in document.tables:
        for row in table.rows:
            row_text = [
                cell.text.strip()
                for cell in row.cells
                if cell.text.strip()
            ]

            if row_text:
                text_parts.append(
                    " | ".join(row_text)
                )

    return "\n\n".join(text_parts)


# ==================================================
# EXTRACT PPTX TEXT
# ==================================================

def extract_pptx_text(
    file_path: Path,
) -> str:
    presentation = Presentation(file_path)

    text_parts = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_parts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]

                    if row_text:
                        slide_parts.append(
                            " | ".join(row_text)
                        )

        if slide_parts:
            text_parts.append(
                f"Slide {slide_number}\n"
                + "\n".join(slide_parts)
            )

    return "\n\n".join(text_parts)


# ==================================================
# EXTRACT DOCUMENT TEXT
# ==================================================

def extract_document_text(
    file_path: Path,
) -> str:
    if not file_path.is_file():
        raise FileNotFoundError(
            "The document file could not be found."
        )

    extension = file_path.suffix.lower()

    if extension not in SUPPORTED_DOCUMENT_TYPES:
        raise ValueError(
            "Text extraction supports only "
            "PDF, DOCX, and PPTX files."
        )

    if extension == ".pdf":
        extracted_text = extract_pdf_text(
            file_path
        )

    elif extension == ".docx":
        extracted_text = extract_docx_text(
            file_path
        )

    else:
        extracted_text = extract_pptx_text(
            file_path
        )

    extracted_text = extracted_text.strip()

    if not extracted_text:
        raise ValueError(
            "No readable text could be extracted "
            "from this document."
        )

    return extracted_text